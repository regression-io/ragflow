import copy
import re
from io import BytesIO
from pptx import Presentation

from rag.app import callback__
from rag.nlp import huqie
from rag.parser.pdf_parser import HuParser


class Ppt(object):
    def __init__(self):
        super().__init__()

    def __extract(self, shape):
        if shape.shape_type == 19:
            tb = shape.table
            rows = []
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
            return "\n".join(rows)

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in shape.shapes:
                t = self.__extract(p)
                if t: texts.append(t)
            return "\n".join(texts)

    def __call__(self, fnm, from_page, to_page, callback=None):
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides[from_page: to_page]):
            texts = []
            for shape in slide.shapes:
                txt = self.__extract(shape)
                if txt: texts.append(txt)
            txts.append("\n".join(texts))
            callback__((i+1)/self.total_page/2, "", callback)

        callback__((min(to_page, self.total_page) - from_page) / self.total_page,
                   "Page {}~{}: Text extraction finished".format(from_page, min(to_page, self.total_page)), callback)
        import aspose.slides as slides
        import aspose.pydrawing as drawing
        imgs = []
        with slides.Presentation(BytesIO(fnm)) as presentation:
            for i, slide in enumerate(presentation.slides[from_page: to_page]):
                buffered = BytesIO()
                slide.get_thumbnail(0.5, 0.5).save(buffered, drawing.imaging.ImageFormat.jpeg)
                imgs.append(buffered.getvalue())
        assert len(imgs) == len(txts), "Slides text and image do not match: {} vs. {}".format(len(imgs), len(txts))
        callback__((min(to_page, self.total_page) - from_page) / self.total_page,
                   "Page {}~{}: Image extraction finished".format(from_page, min(to_page, self.total_page)), callback)

        return [(txts[i], imgs[i]) for i in range(len(txts))]


class Pdf(HuParser):
    def __init__(self):
        super().__init__()

    def __garbage(self, txt):
        txt = txt.lower().strip()
        if re.match(r"[0-9\.,%/-]+$", txt): return True
        if len(txt) < 3:return True
        return False

    def __call__(self, filename, binary=None, from_page=0, to_page=100000, zoomin=3, callback=None):
        self.__images__(filename if not binary else binary, zoomin, from_page, to_page)
        callback__((min(to_page, self.total_page)-from_page) / self.total_page, "Page {}~{}: OCR finished".format(from_page, min(to_page, self.total_page)), callback)
        assert len(self.boxes) == len(self.page_images), "{} vs. {}".format(len(self.boxes), len(self.page_images))
        res = []
        #################### More precisely ###################
        # self._layouts_paddle(zoomin)
        # self._text_merge()
        # pages = {}
        # for b in self.boxes:
        #     if self.__garbage(b["text"]):continue
        #     if b["page_number"] not in pages: pages[b["page_number"]] = []
        #     pages[b["page_number"]].append(b["text"])
        # for i, lines in pages.items():
        #     res.append(("\n".join(lines), self.page_images[i-1]))
        # return res
        ########################################

        for i in range(len(self.boxes)):
            lines = "\n".join([b["text"] for b in self.boxes[i] if not self.__garbage(b["text"])])
            res.append((lines, self.page_images[i]))
        return res


def chunk(filename, binary=None,  from_page=0, to_page=100000, callback=None):
    doc = {
        "docnm_kwd": filename,
        "title_tks": huqie.qie(re.sub(r"\.[a-zA-Z]+$", "", filename))
    }
    doc["title_sm_tks"] = huqie.qieqie(doc["title_tks"])
    res = []
    if re.search(r"\.pptx?$", filename, re.IGNORECASE):
        for txt,img in Ppt()(filename if not binary else binary, from_page, to_page, callback):
            d = copy.deepcopy(doc)
            d["content_ltks"] = huqie.qie(txt)
            d["content_sm_ltks"] = huqie.qieqie(d["content_ltks"])
            d["image"] = img
            res.append(d)
        return res
    if re.search(r"\.pdf$", filename, re.IGNORECASE):
        for txt,img in Pdf()(filename if not binary else binary, from_page=from_page, to_page=to_page, callback=callback):
            d = copy.deepcopy(doc)
            d["content_ltks"] = huqie.qie(txt)
            d["content_sm_ltks"] = huqie.qieqie(d["content_ltks"])
            d["image"] = img
            res.append(d)
        return res
    callback__(-1, "This kind of presentation document did not support yet!", callback)


if __name__== "__main__":
    import sys
    print(chunk(sys.argv[1]))
